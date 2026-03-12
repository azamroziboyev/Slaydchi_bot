import os
import random
from pptx import Presentation
import json
from dotenv import load_dotenv
import requests
import time
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATES_DIR = r"C:\Users\arozi\OneDrive\Desktop\Bestu\slide\templates"

TOPIC_MAP = {
    "agriculture": os.path.join(TEMPLATES_DIR, "agriculture"),
    "farm": os.path.join(TEMPLATES_DIR, "agriculture"),
    "soil": os.path.join(TEMPLATES_DIR, "agriculture"),

    "education": os.path.join(TEMPLATES_DIR, "education"),
    "teacher": os.path.join(TEMPLATES_DIR, "education"),
    "school": os.path.join(TEMPLATES_DIR, "education"),
}

DEFAULT_FOLDER = os.path.join(TEMPLATES_DIR, "default")
load_dotenv()
DEEPSEEK_KEY = os.getenv("SLAYD_KEY")

def expand_grok(prompt, language):
    response = requests.post(
    "https://openrouter.ai/api/v1/chat/completions",
    headers={
        "Authorization": f"Bearer {DEEPSEEK_KEY}",
        "Content-Type": "application/json",
    },
    json={
        "model": "mistralai/devstral-2512:free",
        "messages": [
            {"role": "user", "content": f"{prompt}"}
        ]
    }
    )
    try:
        result = response.json()
        return result["choices"][0]["message"]["content"]

    except Exception as e:
        print("OpenRouter error:", e)
        return "error"
    


def choose_template(topic: str) -> str:
    topic = topic.lower()

    selected_folder = None
    for keyword, folder in TOPIC_MAP.items():
        if keyword in topic:
            selected_folder = folder
            break

    if not selected_folder:
        selected_folder = DEFAULT_FOLDER

    if not os.path.exists(selected_folder):
        selected_folder = DEFAULT_FOLDER

    ppt_files = [f for f in os.listdir(selected_folder) if f.endswith(".pptx")]

    if not ppt_files:
        return os.path.join(DEFAULT_FOLDER, "default1.pptx")

    #return os.path.join(selected_folder, random.choice(ppt_files))
    return os.path.join(selected_folder, 'abs2.pptx')

json_path = r"C:\Users\arozi\OneDrive\Desktop\Bestu\slide\templates\agriculture\indexes.json"
with open(json_path, "r", encoding="utf-8") as f:
        mappings = json.load(f)

style_path = r"C:\Users\arozi\OneDrive\Desktop\Bestu\slide\styles.json"
with open(style_path, "r", encoding="utf-8") as d:
        styles = json.load(d)

# -------------------- APPLY STYLE (SOLE STYLING LOGIC) --------------------        
def apply_style(shape, file_key, slide_index, text_shape_index):
    """
    Applies style to a text shape based on styles.json
    """
    if not shape.has_text_frame:
        return

    file_styles = styles.get(file_key)
    if file_styles is None:
        return

    # use 'is None' so we can distinguish "no slide entry" from "empty dict"
    slide_styles = file_styles.get(str(slide_index))
    if slide_styles is None:
        return

    style_entry = slide_styles.get(str(text_shape_index))
    if style_entry is None:
        # helpful debug: slide has style entries but this particular shape isn't configured
        if slide_styles:
            print(f"⚠ No style entry for slide {slide_index} shape {text_shape_index} in {file_key}")
        return

    # -------- Detect font name & size --------
    font_name = style_entry.get("font_name")
    font_size = style_entry.get("font_size")

    # Compact JSON syntax: {"Font Name": 32}
    if not font_name and not font_size:
        for key, value in style_entry.items():
            if key not in (
                "bold", "italic", "color",
                "alignment", "align",
                "font_name", "font_size"
            ):
                font_name = key
                font_size = value
                break

    bold = style_entry.get("bold")
    italic = style_entry.get("italic")
    color = style_entry.get("color")
    alignment = style_entry.get("alignment") or style_entry.get("align")

    # -------- Apply to all paragraphs --------
    for paragraph in shape.text_frame.paragraphs:

        # Alignment
        if alignment:
            paragraph.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }.get(str(alignment).lower(), PP_ALIGN.LEFT)

        # Reset runs (important!)
        text = paragraph.text
        paragraph.clear()

        run = paragraph.add_run()
        run.text = text
        font = run.font

        if font_name:
            font.name = font_name

        if font_size:
            try:
                font.size = Pt(float(font_size))
            except Exception:
                pass

        if bold is not None:
            font.bold = bool(bold)

        if italic is not None:
            font.italic = bool(italic)

        if color:
            try:
                if isinstance(color, list) and len(color) == 3:
                    r, g, b = color
                elif isinstance(color, str):
                    h = color.lstrip("#")
                    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                else:
                    continue
                font.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass


# -------------------- PRESENTATION WALKER --------------------

def style_presentation(doc, output_pptx, source_pptx=None):
    """
    Walk every slide and every shape index, apply style where configured.
    Accepts a Presentation object directly (no need to read from disk).
    source_pptx: optional path to the original template file used for mapping and styles lookup.
    """
    # Prefer the original template filename for mappings and styles lookups
    file_key = os.path.basename(source_pptx or output_pptx)

    if file_key not in mappings:
        print(f"❌ Warning: '{file_key}' not found in JSON. Styles may not apply.")

    if file_key not in styles:
        print(f"⚠ Warning: No style entry for '{file_key}' in styles.json. No styling will be applied.")

    for slide_index, slide in enumerate(doc.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            apply_style(shape, file_key, slide_index, shape_index)

    doc.save(output_pptx)



def add_titul(doc, topic, student,):  
    """Fill the title slide (first slide) with topic and student name.

    Verifies that the template is present in mappings and writes into the configured text shapes.
    """
    prs = Presentation(doc)

    # The JSON uses:  "organic.pptx": { ... }
    file_key = os.path.basename(doc)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return

    slides_mapping = mappings[file_key]
    titul_number = slides_mapping["0"]["title"]
    subtitle_number = slides_mapping["0"]["matn"]

    prs = Presentation(doc)
    slide = prs.slides[0]     # slide number


    # Fill shapes with the provided text
    slide.shapes[titul_number].text_frame.text = topic
    #apply_style(slide.shapes[titul_number], file_key, 0, titul_number)

    slide.shapes[subtitle_number].text_frame.text = student
    #apply_style(slide.shapes[subtitle_number], file_key, 0, subtitle_number)
    print("Titul list qo'shildi")
    return prs


def add_reja(doc, plan_items, file_path):
    """Populate the second slide with the REJA (plan) entries."""
    #prs = Presentation(doc)
    file_key = os.path.basename(file_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return

    slides_mapping = mappings[file_key]  #bular indexlarni ifodalaydi
    reja_number = slides_mapping["1"]["reja"]
    reja1_number = slides_mapping["1"]["reja1"]
    reja2_number = slides_mapping["1"]["reja2"]
    reja3_number = slides_mapping["1"]["reja3"]

    #prs = Presentation(doc)
    slide = doc.slides[1]     # slide number
    #print(plan_items)
    reja1 = plan_items[0]
    reja2 = plan_items[1]
    reja3 = plan_items[2]
        

    slide.shapes[reja_number].text_frame.text = "Reja:"
    #apply_style(slide.shapes[reja_number], file_key, 1, reja_number)

    slide.shapes[reja1_number].text_frame.text = reja1
    #apply_style(slide.shapes[reja1_number], file_key, 1, reja1_number)

    slide.shapes[reja2_number].text_frame.text = reja2
    #apply_style(slide.shapes[reja2_number], file_key, 1, reja2_number)

    slide.shapes[reja3_number].text_frame.text = reja3
    #apply_style(slide.shapes[reja3_number], file_key, 1, reja3_number)
    print("reja  qo'shildi")
    return doc




# def add_reja(doc, plan_items, file_path):
#     #prs = Presentation(doc)
#     file_key = os.path.basename(file_path)
#     if file_key not in mappings:
#         print(f"❌ Error: '{file_key}' not found in JSON.")
#         return

#     slides_mapping = mappings[file_key]  #bular indexlarni ifodalaydi
#     reja_number = slides_mapping["1"]["reja"]
#     reja1_number = slides_mapping["1"]["reja1"]
#     reja2_number = slides_mapping["1"]["reja2"]
#     reja3_number = slides_mapping["1"]["reja3"]

#     #prs = Presentation(doc)
#     slide = doc.slides[1]     # slide number
    
#     reja1 = plan_items[0]
#     reja2 = plan_items[1]
#     reja3 = plan_items[2]
        

#     slide.shapes[reja_number].text_frame.text = "Reja:"
#     slide.shapes[reja1_number].text_frame.text = reja1
#     slide.shapes[reja2_number].text_frame.text = reja2
#     slide.shapes[reja3_number].text_frame.text = reja3
#     print("reja  qo'shildi")
#     return doc
    
    


def add_content(doc, file_path, topic, plan_items, language):
    """Populate slides with AI-generated content according to mapping JSON.

    For each configured slide (skipping title/reja slides) this will fill configured text shapes.
    """

    file_key = os.path.basename(file_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return
    
    slides_mapping = mappings[file_key]

    # Iterate over slide indexes and their shape mappings
    for slide_index_str, shape_dict in slides_mapping.items():
        
        slide_index = int(slide_index_str)
        if slide_index >= len(doc.slides):
            print(f"⚠ Slide {slide_index} does not exist.")
            continue

        # skip the first and second slides
        if slide_index < 2:
            continue

        if slide_index > 10:
            break

        slide = doc.slides[slide_index]

        for key, shape_index in shape_dict.items():

            if shape_index >= len(slide.shapes):
                print(f"⚠ Shape {shape_index} missing on slide {slide_index}")
                continue

            shape = slide.shapes[shape_index]

            if not shape.has_text_frame:
                print(f"⚠ Shape {shape_index} has no text frame.")
                continue
            
            ai_prompt = (
                f"Write an academic text in Uzbek for a PowerPoint slide.\n"
                f"▪ Topic: '{topic}'\n"
                f"▪ Style: clear, academic, student-level.\n"
                f"▪ total words: 50. (but do not write word counts at the end.)\n"
                f"▪ No wrods counts\n"
            )
            
            final_text = expand_grok(ai_prompt, language)
            final_text = final_text.replace("*", "").replace("#", "").replace("-", "").replace("**", "").replace("##", "").replace("###", "")
            print(key)
            
            # For REJA placeholders, use plan items; otherwise insert AI text
            if key == "reja1":
                shape.text = plan_items[0]
                #apply_style(shape, file_key, slide_index, shape_index)
            elif key == "reja2":
                shape.text = plan_items[1]
                #apply_style(shape, file_key, slide_index, shape_index)
            elif key == "reja3":
                shape.text = plan_items[2]
                #apply_style(shape, file_key, slide_index, shape_index)
            else:
                shape.text = final_text
                #apply_style(shape, file_key, slide_index, shape_index)

            print(f"✅ Slide {slide_index}, Shape {shape_index} filled with '{key}'.")

            time.sleep(0.5) 

    return doc


def add_xulosa(doc, plan_items, file_path, language):

    
    file_key = os.path.basename(file_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return
    
    slides_mapping = mappings[file_key]

    for slide_index_str, shape_dict in slides_mapping.items():
        
        slide_index = int(slide_index_str)
        if slide_index >= len(doc.slides):
            print(f"⚠ Slide {slide_index} does not exist.")
            continue

        #skpi the first and second slides
        if not slide_index == 8:
            continue
        
        if not "conclusion" in shape_dict.keys():
            continue
        # if slide_index > 12:
        #     break

        slide = doc.slides[slide_index]

        for key, shape_index in shape_dict.items():

            if shape_index >= len(slide.shapes):
                print(f"⚠ Shape {shape_index} missing on slide {slide_index}")
                continue

            shape = slide.shapes[shape_index]

            if not shape.has_text_frame:
                print(f"⚠ Shape {shape_index} has no text frame.")
                continue

            if language == "uz":
                xulosa_text = "Xulosa"
            elif language == "ru":
                xulosa_text = "заключение"
            elif language == "eng":
                xulosa_text = "Conclusion"
            else:
                xulosa_text = "Goodbye"

            if key == "conclusion":
                shape.text == xulosa_text
            
            if key == "matn":
                prompt_conclusion = f""" Generate a conclusion about these texts you wrote here {plan_items}. it must be 80 words. texts only."""
                concluded = expand_grok(prompt_conclusion, language)
                concluded = concluded.replace("*", "").replace("#", "").replace("-", "")
                
                #xuloasini yozish
                shape.text = concluded
                #apply_style(shape, file_key, slide_index, shape_index)
            time.sleep(0.5)
    return doc

def add_adabiyotlar(doc, file_path, language, plan_text):
    
    file_key = os.path.basename(file_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return
    
    slides_mapping = mappings[file_key]

    for slide_index_str, shape_dict in slides_mapping.items():
        
        slide_index = int(slide_index_str)
        if slide_index >= len(doc.slides):
            print(f"⚠ Slide {slide_index} does not exist.")
            continue

        #skpi the first and second slides
        if not slide_index == 9:
            continue
        
        if not "references" in shape_dict.keys():
            continue
        # if slide_index > 12:
        #     break

        slide = doc.slides[slide_index]

        for key, shape_index in shape_dict.items():

            if shape_index >= len(slide.shapes):
                print(f"⚠ Shape {shape_index} missing on slide {slide_index}")
                continue

            shape = slide.shapes[shape_index]

            if not shape.has_text_frame:
                print(f"⚠ Shape {shape_index} has no text frame.")
                continue
            if language == "uz":
                adabiyot_text = "Adabiyotlar"
            elif language == "ru":
                adabiyot_text = "литература"
            elif language == "eng":
                adabiyot_text = "References"
            else:
                adabiyot_text = "Goodbye"

            if key == "conclusion":
                shape.text == adabiyot_text

            if key == "matn":
                prompt_reference = f"""Generate a list of 5 references that you have used to describe these points {plan_text}. -output only 5 references no extra text. - important! user also wanted it to be in {language} language"""
                refered = expand_grok(prompt_reference, language)
                refered = refered.replace("*", "").replace("#", "").replace("-", "")
                #adabiyotlar royxatini yozish
                shape.text = refered
            #apply_style(shape, file_key, slide_index, shape_index)

            time.sleep(0.5)
    return doc

def add_farewel(doc, language, file_path):
    file_key = os.path.basename(file_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return

    slides_mapping = mappings[file_key]
    last_key = list(slides_mapping.keys())[-1]
    last_item = slides_mapping[last_key]
    inner_key = list(last_item.keys())[0]
    shape_index = last_item[inner_key]
    slide_index = int(last_key)

    # get slide and shape objects (was using 'slide' before it was defined)
    slide = doc.slides[slide_index]
    shape = slide.shapes[shape_index]

    if language == "uz":
        farewell_text = "Etiboringiz uchun rahmat"
    elif language == "ru":
        farewell_text = "Спасибо за внимание"
    elif language == "eng":
        farewell_text = "Thank you for your attention"
    else:
        farewell_text = "Goodbye"

    if inner_key == "farewel":
        # set text then apply style
        shape.text = farewell_text
        #apply_style(shape, file_key, slide_index, shape_index)

    return doc


def run_engine(topic, student, language, file_path):
    if expand_grok("hi", "eng") == "error":
        return "API"
    else:
        make_presentation(topic, student, language, file_path)


def make_presentation(topic, student, language, final_path):
    """Full slide generation + styling pipeline"""

    #topic = "Farming in Uzbekistan lately"
    #student = "Azamjon Roziboyev"
    #language = "uz"

    template_path = choose_template(topic)

    # ---------- BUILD CONTENT ----------
    doc = add_titul(template_path, topic, student)

    plan_prompt = f"""
    Generate a structured plan (REJA) with maximum 3 points for a referat on '{topic}'.
    - Output only 3 concise points.
    - important! user also wanted it to be in {language} language.
    - No Kirish.
    """

    # if expand_grok("hi", "eng") == "error":
    #     raise RuntimeError("API da xatolik")

    plan_text = expand_grok(plan_prompt, language)

    plan_text = (
        plan_text
        .replace("*", "")
        .replace("#", "")
        .replace("-", "")
        .replace("**", "")
        .replace("##", "")
        .replace("###", "")
    )

    plan_items = [
        line.strip()
        for line in plan_text.split("\n")
        if line.strip()
    ][:3]

    doc = add_reja(doc, plan_items, template_path)
    doc = add_content(doc, template_path, topic, plan_items, language)
    doc = add_xulosa(doc, plan_items, template_path, language)
    doc = add_adabiyotlar(doc, template_path, language, plan_text)
    doc = add_farewel(doc, language, template_path)

    # ---------- SAVE RAW ----------
    #final_path = r"C:\Users\arozi\Desktop\Bestu\slide\abs1_filled_styled.pptx"
    # Pass the original template path so mappings and styles lookup use the template filename
    style_presentation(doc, final_path, template_path)

    print("✅ Presentation generated and styled successfully")

    #return final_path

def main():
    pass

if __name__ == "__main__":
    main()