import os
import json
import time
from pptx import Presentation
from dotenv import load_dotenv
import requests

# -----------------------------  -
# 1. LOAD API KEYS
# ------------------------------
load_dotenv()
DEEPSEEK_KEY = os.getenv("DEEPSEEK_API_KEY")
GROK_KEY = os.getenv("GROK_API_KEY")

# ------------------------------
# 2. AI TEXT GENERATOR (DeepSeek)
# ------------------------------
def generate_deepseek(prompt):
    """Generate a short academic paragraph using Deepseek API.

    Args:
        prompt: prompt string describing the desired content.

    Returns:
        Generated text string or a simple fallback/error message.
    """
    if not DEEPSEEK_KEY:
        return "[AI] Qisqacha matn shu yerda bo'ladi."
    
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "deepseek/deepseek-chat",
        "messages": [
            {"role": "system", "content": "You write clear academic texts."},
            {"role": "user", "content": prompt}
        ]
    }

    try:
        response = requests.post(url, headers=headers, json=data, timeout=15)
        result = response.json()
        return result["choices"][0]["message"]["content"]
    except:
        return "[Error generating paragraph — remote API unavailable]"

# ------------------------------
# 3. AI EXPANDER (Grok)
# ------------------------------
def expand_grok(paragraph, target_words=80):
    """Expand a short paragraph to approximately target_words using Grok API.

    Returns the expanded text (cleaned of common markdown characters) or the original paragraph on failure.
    """
    if not GROK_KEY:
        return paragraph
    
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {GROK_KEY}",
        "Content-Type": "application/json"
    }

    data = {
        "model": "x-ai/grok-4.1-fast:free",
        "messages": [
            {"role": "user", "content": 
                f"Expand this academic Uzbek paragraph to ~{target_words} words. "
                f"No suggestions:\n\n{paragraph}"
            }
        ]
    }

    try:
        response = requests.post(url, headers=headers, json=data, timeout=15)
        result = response.json()
        text = result["choices"][0]["message"]["content"]
        return text.replace("*", "").replace("#", "").replace("-", "")
    except:
        return paragraph
# ------------------------------
# 4. FILL PPTX
# ------------------------------
def add_titul(topic, pptx_path):
    """Fill the first slide's title shape with the given topic and save a filled PPTX copy."""
    prs = Presentation(pptx_path)
    slide = prs.slides[0]     # slide number

    slide.shapes[5].text_frame.text = topic
    output_file = pptx_path.replace(".pptx", "_filled.pptx")
    prs.save(output_file)
    print(f"🎉 Saved: {output_file}")





# ------------------------------
# 4. FILL PPTX
# ------------------------------
def fill_pptx_with_ai(pptx_path, json_path, topic):
    """Fills shapes in the PPTX according to JSON structure & topic."""
    
    # Load JSON structure
    with open(json_path, "r", encoding="utf-8") as f:
        mappings = json.load(f)

    prs = Presentation(pptx_path)

    # The JSON uses:  "organic.pptx": { ... }
    file_key = os.path.basename(pptx_path)
    if file_key not in mappings:
        print(f"❌ Error: '{file_key}' not found in JSON.")
        return

    slides_mapping = mappings[file_key]

    # Process each slide
    for slide_index_str, shape_dict in slides_mapping.items():
        
        slide_index = int(slide_index_str)
        if slide_index >= len(prs.slides):
            print(f"⚠ Slide {slide_index} does not exist.")
            continue
        #skip the first and second slides
        if slide_index < 3:
            continue

        slide = prs.slides[slide_index]


        # each key like: "title": 5, "matn": 22
        for key, shape_index in shape_dict.items():

            if shape_index >= len(slide.shapes):
                print(f"⚠ Shape {shape_index} missing on slide {slide_index}")
                continue

            shape = slide.shapes[shape_index]
            if not shape.has_text_frame:
                print(f"⚠ Shape {shape_index} has no text frame.")
                continue

            # -------------------------------
            # AI PROMPT BASED ON TOPIC
            # -------------------------------
            ai_prompt = (
                f"Write an academic text in Uzbek for a PowerPoint slide.\n"
                f"▪ Topic: '{topic}'\n"
                f"▪ Section: '{key}'\n"
                f"▪ Style: clear, academic, student-level.\n"
                f"▪ No bullet points.\n"
            )

            short_text = generate_deepseek(ai_prompt)
            final_text = expand_grok(short_text, target_words=60)

            shape.text = final_text

            print(f"✅ Slide {slide_index}, Shape {shape_index} filled with '{key}'.")

            time.sleep(0.5)

    output_file = pptx_path.replace(".pptx", "_filled.pptx")
    prs.save(output_file)
    print(f"🎉 Saved: {output_file}")

# ------------------------------
# 5. MAIN
# ------------------------------
def main():
    """Demo runner that fills the title slide for a sample PPTX using the module functions."""
    PPTX_FILE = r"C:\Users\arozi\Desktop\Bestu\slide\templates\agriculture\organic.pptx"
    JSON_FILE = r"C:\Users\arozi\Desktop\Bestu\slide\templates\agriculture\indexes.json"
    TOPIC = "Qishloq xo‘jaligida organik o‘g‘itlarning ahamiyati"

    #fill_pptx_with_ai(PPTX_FILE, JSON_FILE, TOPIC)
    add_titul(TOPIC, PPTX_FILE)

if __name__ == "__main__":
    main()
