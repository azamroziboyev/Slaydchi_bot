import os
from pptx import Presentation

FOLDER = r"C:\Users\arozi\OneDrive\Desktop\Bestu\slide\templates\agriculture"

for filename in os.listdir(FOLDER):
    if filename.endswith(".pptx"):
        pptx_path = os.path.join(FOLDER, filename)
        print(f"\nTrying to open: {pptx_path}")

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            print(f"❌ Skipped (invalid or corrupted file): {filename}")
            print(f"   Reason: {e}")
            continue

        # ---------- Process valid PPTX ----------
        for slide_number, slide in enumerate(prs.slides):
            for index, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    shape.text = f"index {index}"

        output_path = pptx_path.replace(".pptx", "_indexed.pptx")
        prs.save(output_path)

        print(f"✔ Saved: {output_path}")
    
    
# Texts to write in the textboxes
# texts = [
#     "index 0",
#     "index 1",
#     "Point 2",
#     "Point 3",
#     "Point 4",
#     "Point 5",
#     "Point 6"
# ]

# # Loop through shapes and write text
# for i, shape in enumerate(slide.shapes):
#     if shape.has_text_frame and i < len(texts):
#         shape.text = texts[i]


# "organic.pptx": {
#         "0": {"title": 5, "matn": 22},
#         "1": {"reja": 8, "reja1": 5, "reja2": 6, "reja3": 7},
#         "2": {"title": 6, "subtitle": 7},
#         "3": {"matn": 4},
#         "4": {"matn": 1},
#         "5": {"matn": 14},
#         "6": {"title": 6, "matn": 7},
#         "7": {"matn": 0},
#         "8": {"title": 6,"matn": 7},
#         "9": {"matn": 0},
#         "10": {"matn": 0},
#         "11": {"matn": 42}
#     },